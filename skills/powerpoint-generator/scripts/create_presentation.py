#!/usr/bin/env python3
"""Create PowerPoint presentations from command-line arguments.

Usage:
    python create_presentation.py --title "Title" --slides "Slide1|Content1" "Slide2|Content2" --output presentation.pptx
"""

import argparse
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def create_presentation(title, subtitle, slides, output_path):
    """Create a PowerPoint presentation.
    
    Args:
        title: Presentation title
        subtitle: Presentation subtitle
        slides: List of (title, content) tuples
        output_path: Path to save the presentation
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle

    # Add content slides
    for slide_title, slide_content in slides:
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = content_slide.shapes.title
        content_shape = content_slide.placeholders[1]
        
        title_shape.text = slide_title
        
        # Handle bullet points
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Split content by newlines for bullet points
        lines = slide_content.split('\n')
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
                
            if i == 0:
                text_frame.text = line
            else:
                p = text_frame.add_paragraph()
                p.text = line
                p.level = 0

    # Save presentation
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    print(f"✓ Slides: {len(slides) + 1} (1 title + {len(slides)} content)")
    
    return output_path


def main():
    parser = argparse.ArgumentParser(description='Create PowerPoint presentations')
    parser.add_argument('--title', required=True, help='Presentation title')
    parser.add_argument('--subtitle', default='', help='Presentation subtitle')
    parser.add_argument('--slides', nargs='+', required=True, 
                       help='Slides in format "Title|Content with\\nBullets"')
    parser.add_argument('--output', default='presentation.pptx', 
                       help='Output filename')
    parser.add_argument('--output-dir', default=None,
                       help='Output directory (default: current directory)')
    
    args = parser.parse_args()
    
    # Parse slides
    slides = []
    for slide_str in args.slides:
        if '|' in slide_str:
            slide_title, slide_content = slide_str.split('|', 1)
            slides.append((slide_title.strip(), slide_content.strip()))
        else:
            slides.append((slide_str.strip(), ''))
    
    # Determine output path
    if args.output_dir:
        output_dir = Path(args.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / args.output
    else:
        output_path = Path(args.output)
    
    # Create presentation
    try:
        create_presentation(args.title, args.subtitle, slides, str(output_path))
        
        # Print file output marker for chat UI
        print(f'FILE_OUTPUT: {{"path": "{output_path.absolute()}", "filename": "{args.output}", "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}}')
        
        return 0
    except Exception as e:
        print(f"Error creating presentation: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())
